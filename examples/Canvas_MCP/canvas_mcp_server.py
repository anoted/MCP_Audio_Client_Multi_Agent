"""Canvas LMS MCP server.

Exposes tools, resources, and prompts for:
  * setting and grading homework (assignments + submissions, with
    description / deadline / points)
  * reading submission text for assessment — online text entries and
    uploaded PDF / DOCX files (text is extracted server-side; the LLM
    client connected over MCP performs the actual assessment)
  * creating quizzes and quiz questions
  * managing course modules, including reading the content behind module
    items (pages, PDF/DOCX/PPTX/text files, assignment & quiz descriptions)
  * wiki pages (list / read / create / update)
  * course files (list, upload from local disk, extract text)
  * announcements, discussions, and the student roster

Transport: stdio (default — the voice client registers the server via
mcp_servers.json). Streamable HTTP is available by swapping the mcp.run()
call at the bottom of this file; the HTTP endpoint has no auth, so bind it
to localhost only. Configuration comes from the environment / .env:

  CANVAS_BASE_URL     e.g. https://yourschool.instructure.com
  CANVAS_API_TOKEN    Canvas personal access token (Account > Settings >
                      New Access Token)
  MCP_HOST/MCP_PORT   HTTP transport only (default 127.0.0.1:8017)
  CANVAS_MCP_AUDIT    set to 0 to disable the server-side audit log
  CANVAS_MCP_LOG_DIR  audit log directory (default: logs/ next to this file)

Every tool call is appended to a JSONL audit log (arguments PII-masked), so
the server keeps its own record independent of the host app's session log.

Run:  python canvas_mcp_server.py
"""

from __future__ import annotations

import io
import json
import mimetypes
import os
from html.parser import HTMLParser
from pathlib import Path

import httpx
from dotenv import load_dotenv
from mcp.server.fastmcp import FastMCP

load_dotenv()

MCP_HOST = os.environ.get("MCP_HOST", "127.0.0.1")
MCP_PORT = int(os.environ.get("MCP_PORT", "8017"))

MAX_ATTACHMENT_BYTES = 25 * 1024 * 1024  # refuse to download files larger than 25 MB

mcp = FastMCP(
    "canvas-lms",
    instructions=(
        "Canvas LMS course management server. Use the tools to list courses, "
        "create/grade homework assignments, read student submission text "
        "(including uploaded PDF and DOCX files), build quizzes, create and "
        "edit wiki pages, upload files (PDF/PPTX/etc.) to the course, post "
        "announcements and discussions, and manage modules — including reading "
        "the content behind any module item with read_module_item. Grades are "
        "only written when grade_submission is called — fetch submission text "
        "with get_submission_text, assess it against the rubric, then post the "
        "grade and feedback comment."
    ),
    host=MCP_HOST,
    port=MCP_PORT,
    streamable_http_path="/mcp",
)


# --------------------------------------------------------------------------
# Canvas REST helpers
# --------------------------------------------------------------------------

_http: httpx.AsyncClient | None = None


def _client() -> httpx.AsyncClient:
    global _http
    if _http is None:
        _http = httpx.AsyncClient(timeout=60.0, follow_redirects=True)
    return _http


def _config() -> tuple[str, str]:
    base = os.environ.get("CANVAS_BASE_URL", "").rstrip("/")
    token = os.environ.get("CANVAS_API_TOKEN", "")
    if not base or not token:
        raise RuntimeError(
            "Canvas is not configured. Set CANVAS_BASE_URL and CANVAS_API_TOKEN "
            "in the environment or a .env file next to canvas_mcp_server.py, "
            "then restart the server."
        )
    return base, token


def _headers() -> dict[str, str]:
    _, token = _config()
    return {"Authorization": f"Bearer {token}"}


# --------------------------------------------------------------------------
# Client-side rate limiting.
#
# Canvas meters each access token with a leaky bucket (~700 units, refilled
# continuously; the X-Rate-Limit-Remaining header reports the level). Every
# request costs units proportional to its server-side compute time, and each
# *concurrent* request additionally pre-holds 50 units, so bursts and
# parallelism drain the bucket fast — when it empties Canvas answers
# 403 "Rate Limit Exceeded". Instead of failing, we:
#   1. cap concurrency (CANVAS_MAX_CONCURRENT, default 3),
#   2. space request starts (CANVAS_MIN_INTERVAL_MS, default 200 ms),
#   3. ease off proactively when the bucket runs low (< ~150 units),
#   4. retry rate-limited responses with exponential backoff + jitter
#      (CANVAS_RATE_RETRIES, default 5), honoring Retry-After.
# --------------------------------------------------------------------------

import asyncio
import random
import sys
import time

RATE_MAX_CONCURRENT = int(os.environ.get("CANVAS_MAX_CONCURRENT", "3"))
RATE_MIN_INTERVAL_S = float(os.environ.get("CANVAS_MIN_INTERVAL_MS", "200")) / 1000
RATE_MAX_RETRIES = int(os.environ.get("CANVAS_RATE_RETRIES", "5"))
_LOW_BUCKET = 150.0  # start slowing down below this X-Rate-Limit-Remaining

_rate_sem: asyncio.Semaphore | None = None
_pace_lock: asyncio.Lock | None = None
_last_start = 0.0


def _rate_limited(resp: httpx.Response) -> bool:
    return resp.status_code == 429 or (
        resp.status_code == 403 and "rate limit exceeded" in resp.text.lower()
    )


async def _throttled(method: str, url: str, **kwargs) -> httpx.Response:
    """All Canvas HTTP goes through here: paced, bucket-aware, retrying."""
    global _rate_sem, _pace_lock, _last_start
    if _rate_sem is None:  # lazily bound to the running event loop
        _rate_sem = asyncio.Semaphore(RATE_MAX_CONCURRENT)
        _pace_lock = asyncio.Lock()
    resp: httpx.Response | None = None
    for attempt in range(RATE_MAX_RETRIES + 1):
        async with _rate_sem:
            async with _pace_lock:  # spread out request starts
                wait = _last_start + RATE_MIN_INTERVAL_S - time.monotonic()
                if wait > 0:
                    await asyncio.sleep(wait)
                _last_start = time.monotonic()
            resp = await _client().request(method, url, **kwargs)
        if not _rate_limited(resp):
            # Proactive brake: the emptier the bucket, the longer we pause.
            try:
                remaining = float(resp.headers.get("x-rate-limit-remaining", ""))
            except ValueError:
                remaining = _LOW_BUCKET
            if remaining < _LOW_BUCKET:
                brake = min(3.0, (_LOW_BUCKET - remaining) / _LOW_BUCKET * 3.0)
                print(f"[canvas] rate bucket low ({remaining:.0f}) — easing off "
                      f"{brake:.1f}s", file=sys.stderr)
                await asyncio.sleep(brake)
            return resp
        if attempt == RATE_MAX_RETRIES:
            break
        try:
            delay = float(resp.headers.get("retry-after", ""))
        except ValueError:
            delay = min(20.0, 2.0 ** attempt + random.random())
        print(f"[canvas] rate limited — retry {attempt + 1}/{RATE_MAX_RETRIES} "
              f"in {delay:.1f}s", file=sys.stderr)
        await asyncio.sleep(delay)
    return resp


async def _canvas(method: str, path: str, *, params=None, body=None):
    """Single Canvas API call. `path` is relative to /api/v1."""
    base, _ = _config()
    resp = await _throttled(
        method,
        f"{base}/api/v1{path}",
        params=params,
        json=body,
        headers=_headers(),
    )
    if resp.status_code >= 400:
        raise RuntimeError(f"Canvas API {resp.status_code} on {method} {path}: {resp.text[:500]}")
    if resp.status_code == 204 or not resp.content:
        return {}
    return resp.json()


async def _canvas_paginated(path: str, *, params=None) -> list:
    """GET all pages of a Canvas list endpoint (follows Link: rel=next)."""
    base, _ = _config()
    url = f"{base}/api/v1{path}"
    merged = {"per_page": 100, **(params or {})}
    out: list = []
    while url:
        resp = await _throttled("GET", url, params=merged, headers=_headers())
        merged = None  # params are baked into the next-page URL
        if resp.status_code >= 400:
            raise RuntimeError(f"Canvas API {resp.status_code} on GET {path}: {resp.text[:500]}")
        data = resp.json()
        out.extend(data if isinstance(data, list) else [data])
        url = resp.links.get("next", {}).get("url")
    return out


# --------------------------------------------------------------------------
# Text extraction (HTML, PDF, DOCX)
# --------------------------------------------------------------------------


class _HTMLText(HTMLParser):
    def __init__(self):
        super().__init__()
        self.parts: list[str] = []

    def handle_data(self, data):
        self.parts.append(data)

    def handle_starttag(self, tag, attrs):
        if tag in ("p", "br", "div", "li", "tr", "h1", "h2", "h3", "h4"):
            self.parts.append("\n")


def _strip_html(html: str) -> str:
    if not html:
        return ""
    p = _HTMLText()
    p.feed(html)
    text = "".join(p.parts)
    return "\n".join(line.strip() for line in text.splitlines() if line.strip())


def _extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    pages = []
    for i, page in enumerate(reader.pages, start=1):
        pages.append(f"[page {i}]\n{(page.extract_text() or '').strip()}")
    return "\n\n".join(pages)


def _extract_docx(data: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(data))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[speaker notes] {notes}")
        slides.append(f"[slide {i}]\n" + "\n".join(parts))
    return "\n\n".join(slides)


def _extract_file_text(name: str, ctype: str, data: bytes) -> str:
    """Extract readable text from file bytes (PDF, DOCX, PPTX, HTML, plain text)."""
    lower = (name or "").lower()
    ctype = (ctype or "").lower()
    if "pdf" in ctype or lower.endswith(".pdf"):
        return _extract_pdf(data)
    if "wordprocessingml" in ctype or lower.endswith(".docx"):
        return _extract_docx(data)
    if "presentationml" in ctype or lower.endswith(".pptx"):
        return _extract_pptx(data)
    if "html" in ctype or lower.endswith((".html", ".htm")):
        return _strip_html(data.decode("utf-8", errors="replace"))
    if ctype.startswith("text/") or lower.endswith(
        (".txt", ".md", ".markdown", ".csv", ".json", ".py", ".r", ".sql")
    ):
        return data.decode("utf-8", errors="replace")
    raise ValueError(
        f"unsupported file type ({ctype or lower}) — PDF, DOCX, PPTX, HTML, and "
        "plain-text files can be extracted (legacy .doc/.ppt and media files cannot)"
    )


async def _file_text(f: dict) -> str:
    """Download a Canvas file record and extract its text."""
    name = f.get("display_name") or f.get("filename") or f"file {f.get('id')}"
    size = f.get("size") or 0
    if size > MAX_ATTACHMENT_BYTES:
        return f"[skipped: file is {size} bytes, over the {MAX_ATTACHMENT_BYTES} byte limit]"
    url = f.get("url")
    if not url:
        return "[file has no download URL — it may be locked or hidden]"
    resp = await _throttled("GET", url, headers=_headers())
    if resp.status_code >= 400:
        return f"[download failed: HTTP {resp.status_code}]"
    try:
        text = _extract_file_text(name, f.get("content-type") or "", resp.content)
    except Exception as exc:
        return f"[text extraction failed: {exc}]"
    return text.strip() or "[no extractable text — file may be scanned images or media]"


async def _attachment_text(att: dict) -> str:
    """Download one submission attachment and extract its text (PDF/DOCX only)."""
    name = att.get("filename") or att.get("display_name") or "file"
    ctype = (att.get("content-type") or att.get("content_type") or "").lower()
    lower = name.lower()
    is_pdf = "pdf" in ctype or lower.endswith(".pdf")
    is_docx = "wordprocessingml" in ctype or lower.endswith(".docx")
    if not (is_pdf or is_docx):
        return f"=== {name} ===\n[skipped: only PDF and DOCX files are supported, got {ctype or lower}]"

    size = att.get("size") or 0
    if size > MAX_ATTACHMENT_BYTES:
        return f"=== {name} ===\n[skipped: file is {size} bytes, over the {MAX_ATTACHMENT_BYTES} byte limit]"

    resp = await _throttled("GET", att["url"], headers=_headers())
    if resp.status_code >= 400:
        return f"=== {name} ===\n[download failed: HTTP {resp.status_code}]"

    try:
        text = _extract_pdf(resp.content) if is_pdf else _extract_docx(resp.content)
    except Exception as exc:  # corrupt / password-protected files etc.
        return f"=== {name} ===\n[text extraction failed: {exc}]"
    return f"=== {name} ===\n{text.strip() or '[no extractable text — file may be scanned images]'}"


# --------------------------------------------------------------------------
# Slim serializers — keep tool output compact for the LLM
# --------------------------------------------------------------------------


def _slim_course(c: dict) -> dict:
    return {
        "id": c.get("id"),
        "name": c.get("name"),
        "course_code": c.get("course_code"),
        "workflow_state": c.get("workflow_state"),
        "term": (c.get("term") or {}).get("name"),
        "total_students": c.get("total_students"),
    }


def _slim_assignment(a: dict) -> dict:
    return {
        "id": a.get("id"),
        "name": a.get("name"),
        "due_at": a.get("due_at"),
        "points_possible": a.get("points_possible"),
        "submission_types": a.get("submission_types"),
        "published": a.get("published"),
        "needs_grading_count": a.get("needs_grading_count"),
        "html_url": a.get("html_url"),
    }


def _slim_submission(s: dict) -> dict:
    user = s.get("user") or {}
    return {
        "user_id": s.get("user_id"),
        "user_name": user.get("name") or user.get("sortable_name"),
        "workflow_state": s.get("workflow_state"),
        "submission_type": s.get("submission_type"),
        "submitted_at": s.get("submitted_at"),
        "late": s.get("late"),
        "score": s.get("score"),
        "grade": s.get("grade"),
        "attachments": [
            {"filename": a.get("filename"), "content_type": a.get("content-type"), "size": a.get("size")}
            for a in (s.get("attachments") or [])
        ],
    }


def _slim_quiz(q: dict) -> dict:
    return {
        "id": q.get("id"),
        "title": q.get("title"),
        "quiz_type": q.get("quiz_type"),
        "question_count": q.get("question_count"),
        "points_possible": q.get("points_possible"),
        "due_at": q.get("due_at"),
        "time_limit": q.get("time_limit"),
        "published": q.get("published"),
        "html_url": q.get("html_url"),
    }


def _slim_module(m: dict) -> dict:
    return {
        "id": m.get("id"),
        "name": m.get("name"),
        "position": m.get("position"),
        "published": m.get("published"),
        "items_count": m.get("items_count"),
        "items": [
            {
                "id": i.get("id"),
                "title": i.get("title"),
                "type": i.get("type"),
                "content_id": i.get("content_id"),
                "position": i.get("position"),
                "published": i.get("published"),
            }
            for i in (m.get("items") or [])
        ],
    }


def _slim_page(p: dict) -> dict:
    return {
        "page_id": p.get("page_id"),
        "url": p.get("url"),
        "title": p.get("title"),
        "published": p.get("published"),
        "front_page": p.get("front_page"),
        "updated_at": p.get("updated_at"),
    }


def _slim_file(f: dict) -> dict:
    return {
        "id": f.get("id"),
        "display_name": f.get("display_name"),
        "content_type": f.get("content-type"),
        "size": f.get("size"),
        "folder_id": f.get("folder_id"),
        "updated_at": f.get("updated_at"),
    }


# --------------------------------------------------------------------------
# Server-side audit log — one JSONL line per tool call.
#
# The host app writes its own session log, but that log comes from the same
# process it governs. This one is the independent witness on the tool side:
# which tool ran, with which arguments (PII-masked, long values truncated),
# whether it succeeded, and how long it took. Logs go to files only — stdout
# carries the MCP protocol on stdio.
# --------------------------------------------------------------------------

import functools
import re

_AUDIT_ENABLED = os.environ.get("CANVAS_MCP_AUDIT", "1").strip().lower() not in (
    "0", "false", "no", "off",
)
_AUDIT_DIR = Path(os.environ.get("CANVAS_MCP_LOG_DIR") or Path(__file__).parent / "logs")
_AUDIT_PATH = _AUDIT_DIR / f"canvas-mcp-{time.strftime('%Y%m%d-%H%M%S')}-{os.getpid()}.jsonl"

# Same masking rules as the host's privacy filter: emails and phone numbers
# never land in a log, whatever argument they arrive in.
_PII_EMAIL = re.compile(r"[\w.+-]+@[\w-]+\.[\w.-]+")
_PII_PHONE = re.compile(
    r"(?<![\d/#-])(?:\+?\d{1,3}[ .-]?)?(?:\(\d{3}\)|\d{3})[ .-]\d{3}[ .-]\d{4}(?!\d)"
)


def _mask(value, limit: int = 300):
    """PII-masked, truncated copy of one tool argument for the audit log."""
    if isinstance(value, str):
        value = _PII_EMAIL.sub("[email]", value)
        value = _PII_PHONE.sub("[phone]", value)
        if len(value) > limit:
            value = value[:limit] + f"…(+{len(value) - limit} chars)"
        return value
    if isinstance(value, dict):
        return {k: _mask(v, limit) for k, v in value.items()}
    if isinstance(value, list):
        return [_mask(v, limit) for v in value[:20]]
    return value


def _audit(entry: dict) -> None:
    if not _AUDIT_ENABLED:
        return
    try:
        _AUDIT_DIR.mkdir(parents=True, exist_ok=True)
        with _AUDIT_PATH.open("a", encoding="utf-8") as fh:
            fh.write(json.dumps(entry, ensure_ascii=False, default=str) + "\n")
    except OSError:
        pass  # never let logging break a tool call


_plain_tool = mcp.tool


def _logged_tool(*dargs, **dkwargs):
    """Drop-in for @mcp.tool() that audits every invocation of the tool."""
    register = _plain_tool(*dargs, **dkwargs)

    def wrap(fn):
        @functools.wraps(fn)
        async def logged(*args, **kwargs):
            t0 = time.monotonic()
            entry = {
                "ts": time.strftime("%Y-%m-%d %H:%M:%S"),
                "tool": fn.__name__,
                "args": _mask(kwargs),
            }
            try:
                result = await fn(*args, **kwargs)
            except Exception as exc:
                entry.update(
                    ok=False,
                    error=_mask(str(exc), 500),
                    ms=round((time.monotonic() - t0) * 1000),
                )
                _audit(entry)
                raise
            entry.update(
                ok=True,
                ms=round((time.monotonic() - t0) * 1000),
                result_chars=len(json.dumps(result, ensure_ascii=False, default=str))
                if result is not None else 0,
            )
            _audit(entry)
            return result

        return register(logged)

    return wrap


mcp.tool = _logged_tool  # every @mcp.tool() below is audited


# --------------------------------------------------------------------------
# Tools — courses & homework (assignments + grading)
# --------------------------------------------------------------------------


@mcp.tool()
async def list_courses(include_concluded: bool = False) -> list[dict]:
    """List the courses this Canvas token can teach/manage.

    Returns id, name, code, term, and student count for each course.
    """
    params: dict = {"include[]": ["term", "total_students"]}
    if not include_concluded:
        params["enrollment_state"] = "active"
    courses = await _canvas_paginated("/courses", params=params)
    return [_slim_course(c) for c in courses]


@mcp.tool()
async def list_assignments(course_id: int) -> list[dict]:
    """List assignments (homework) in a course, including how many submissions still need grading."""
    assignments = await _canvas_paginated(f"/courses/{course_id}/assignments")
    return [_slim_assignment(a) for a in assignments]


@mcp.tool()
async def create_assignment(
    course_id: int,
    name: str,
    description_html: str,
    points_possible: float = 100,
    due_at: str | None = None,
    unlock_at: str | None = None,
    lock_at: str | None = None,
    submission_types: list[str] | None = None,
    allowed_extensions: list[str] | None = None,
    published: bool = False,
    module_id: int | None = None,
) -> dict:
    """Create a homework assignment in a course.

    - description_html: the assignment instructions (HTML allowed).
    - due_at: deadline, ISO 8601, e.g. "2026-07-15T23:59:00-05:00".
    - unlock_at / lock_at: optional window when students can submit.
    - submission_types: e.g. ["online_text_entry", "online_upload"] (default).
    - allowed_extensions: e.g. ["pdf", "docx"] (only applies to online_upload).
    - published: leave False to create as a draft the instructor can review first.
    - module_id: optionally add the new assignment straight to that module.
    """
    assignment: dict = {
        "name": name,
        "description": description_html,
        "points_possible": points_possible,
        "submission_types": submission_types or ["online_text_entry", "online_upload"],
        "published": published,
    }
    if due_at:
        assignment["due_at"] = due_at
    if unlock_at:
        assignment["unlock_at"] = unlock_at
    if lock_at:
        assignment["lock_at"] = lock_at
    if allowed_extensions:
        assignment["allowed_extensions"] = allowed_extensions
    created = await _canvas("POST", f"/courses/{course_id}/assignments", body={"assignment": assignment})
    out = _slim_assignment(created)
    if module_id is not None and created.get("id"):
        out["module_item"] = await add_module_item(
            course_id=course_id, module_id=module_id, item_type="Assignment",
            title=name, content_id=created["id"],
        )
    return out


@mcp.tool()
async def update_assignment(
    course_id: int,
    assignment_id: int,
    name: str | None = None,
    description_html: str | None = None,
    points_possible: float | None = None,
    due_at: str | None = None,
    unlock_at: str | None = None,
    lock_at: str | None = None,
    published: bool | None = None,
) -> dict:
    """Update an existing assignment. Only the fields you pass are changed.

    Set published=True to publish a draft assignment so students can see it.
    """
    assignment: dict = {}
    if name is not None:
        assignment["name"] = name
    if description_html is not None:
        assignment["description"] = description_html
    if points_possible is not None:
        assignment["points_possible"] = points_possible
    if due_at is not None:
        assignment["due_at"] = due_at
    if unlock_at is not None:
        assignment["unlock_at"] = unlock_at
    if lock_at is not None:
        assignment["lock_at"] = lock_at
    if published is not None:
        assignment["published"] = published
    if not assignment:
        raise ValueError("Pass at least one field to update.")
    updated = await _canvas(
        "PUT", f"/courses/{course_id}/assignments/{assignment_id}", body={"assignment": assignment}
    )
    return _slim_assignment(updated)


@mcp.tool()
async def list_submissions(course_id: int, assignment_id: int, only_ungraded: bool = False) -> list[dict]:
    """List student submissions for an assignment (who submitted, when, current grade, attached files).

    Use only_ungraded=True to see just the submissions that still need a grade.
    """
    subs = await _canvas_paginated(
        f"/courses/{course_id}/assignments/{assignment_id}/submissions",
        params={"include[]": ["user"]},
    )
    slim = [_slim_submission(s) for s in subs if s.get("workflow_state") != "unsubmitted"]
    if only_ungraded:
        slim = [s for s in slim if s["workflow_state"] != "graded" or s["grade"] is None]
    return slim


@mcp.tool()
async def get_submission_text(course_id: int, assignment_id: int, user_id: int) -> str:
    """Fetch the full text of one student's submission so it can be assessed.

    Returns the online text entry (HTML stripped) and/or the extracted text of
    every attached PDF or DOCX file. Other file types are listed but skipped.
    Use this before grading: read the text, assess it against the rubric, then
    call grade_submission with the grade and a feedback comment.
    """
    sub = await _canvas(
        "GET",
        f"/courses/{course_id}/assignments/{assignment_id}/submissions/{user_id}",
        params={"include[]": ["user", "submission_comments"]},
    )
    user = sub.get("user") or {}
    header = (
        f"Student: {user.get('name', user_id)} (user_id={user_id})\n"
        f"Submitted: {sub.get('submitted_at')}  Type: {sub.get('submission_type')}  "
        f"Late: {sub.get('late')}  Current grade: {sub.get('grade')}"
    )
    sections = [header]

    if sub.get("body"):
        sections.append("=== Text entry ===\n" + _strip_html(sub["body"]))

    for att in sub.get("attachments") or []:
        sections.append(await _attachment_text(att))

    if len(sections) == 1:
        sections.append("[no submission content found — the student may not have submitted]")
    return "\n\n".join(sections)


@mcp.tool()
async def grade_submission(
    course_id: int,
    assignment_id: int,
    user_id: int,
    grade: str,
    comment: str | None = None,
) -> dict:
    """Post a grade (and optional feedback comment) for one student's submission.

    grade accepts points ("8.5"), a percentage ("85%"), or a letter grade ("B+"),
    matching the assignment's grading type. This writes to the Canvas gradebook.
    """
    body: dict = {"submission": {"posted_grade": grade}}
    if comment:
        body["comment"] = {"text_comment": comment}
    result = await _canvas(
        "PUT", f"/courses/{course_id}/assignments/{assignment_id}/submissions/{user_id}", body=body
    )
    return {
        "user_id": result.get("user_id"),
        "grade": result.get("grade"),
        "score": result.get("score"),
        "workflow_state": result.get("workflow_state"),
        "graded_at": result.get("graded_at"),
    }


# --------------------------------------------------------------------------
# Tools — quizzes
# --------------------------------------------------------------------------


@mcp.tool()
async def list_quizzes(course_id: int) -> list[dict]:
    """List quizzes in a course."""
    quizzes = await _canvas_paginated(f"/courses/{course_id}/quizzes")
    return [_slim_quiz(q) for q in quizzes]


@mcp.tool()
async def create_quiz(
    course_id: int,
    title: str,
    description_html: str = "",
    quiz_type: str = "assignment",
    time_limit_minutes: int | None = None,
    due_at: str | None = None,
    shuffle_answers: bool = True,
    allowed_attempts: int = 1,
    module_id: int | None = None,
) -> dict:
    """Create a quiz (unpublished draft). Add questions with add_quiz_question, then publish_quiz.

    - quiz_type: "assignment" (graded), "practice_quiz", "graded_survey", or "survey".
    - due_at: ISO 8601 timestamp.
    - module_id: optionally add the new quiz straight to that module.
    """
    quiz: dict = {
        "title": title,
        "description": description_html,
        "quiz_type": quiz_type,
        "shuffle_answers": shuffle_answers,
        "allowed_attempts": allowed_attempts,
        "published": False,
    }
    if time_limit_minutes:
        quiz["time_limit"] = time_limit_minutes
    if due_at:
        quiz["due_at"] = due_at
    created = await _canvas("POST", f"/courses/{course_id}/quizzes", body={"quiz": quiz})
    out = _slim_quiz(created)
    if module_id is not None and created.get("id"):
        out["module_item"] = await add_module_item(
            course_id=course_id, module_id=module_id, item_type="Quiz",
            title=title, content_id=created["id"],
        )
    return out


@mcp.tool()
async def add_quiz_question(
    course_id: int,
    quiz_id: int,
    question_text: str,
    question_type: str = "multiple_choice_question",
    points_possible: float = 1,
    answers: list[dict] | None = None,
    question_name: str | None = None,
) -> dict:
    """Add one question to a quiz.

    - question_type: "multiple_choice_question", "true_false_question",
      "short_answer_question", "essay_question", "multiple_answers_question",
      "numerical_question", or "text_only_question".
    - answers: list of {"text": "...", "correct": true/false} objects — required for
      multiple choice / true-false / multiple answers / short answer types.
      Exactly the correct ones get weight 100. Not needed for essay questions.
    """
    canvas_answers = [
        {"answer_text": a["text"], "answer_weight": 100 if a.get("correct") else 0}
        for a in (answers or [])
    ]
    question: dict = {
        "question_name": question_name or question_text[:50],
        "question_text": question_text,
        "question_type": question_type,
        "points_possible": points_possible,
    }
    if canvas_answers:
        question["answers"] = canvas_answers
    created = await _canvas(
        "POST", f"/courses/{course_id}/quizzes/{quiz_id}/questions", body={"question": question}
    )
    return {
        "id": created.get("id"),
        "question_name": created.get("question_name"),
        "question_type": created.get("question_type"),
        "points_possible": created.get("points_possible"),
    }


@mcp.tool()
async def update_quiz(
    course_id: int,
    quiz_id: int,
    title: str | None = None,
    description_html: str | None = None,
    due_at: str | None = None,
    time_limit_minutes: int | None = None,
    allowed_attempts: int | None = None,
    published: bool | None = None,
) -> dict:
    """Update quiz settings (title, description, due date, time limit, attempts, publish state).

    Only the fields you pass are changed.
    """
    quiz: dict = {}
    if title is not None:
        quiz["title"] = title
    if description_html is not None:
        quiz["description"] = description_html
    if due_at is not None:
        quiz["due_at"] = due_at
    if time_limit_minutes is not None:
        quiz["time_limit"] = time_limit_minutes
    if allowed_attempts is not None:
        quiz["allowed_attempts"] = allowed_attempts
    if published is not None:
        quiz["published"] = published
    if not quiz:
        raise ValueError("Pass at least one field to update.")
    updated = await _canvas("PUT", f"/courses/{course_id}/quizzes/{quiz_id}", body={"quiz": quiz})
    return _slim_quiz(updated)


@mcp.tool()
async def publish_quiz(course_id: int, quiz_id: int) -> dict:
    """Publish a quiz so students can take it. Do this after all questions are added."""
    updated = await _canvas("PUT", f"/courses/{course_id}/quizzes/{quiz_id}", body={"quiz": {"published": True}})
    return _slim_quiz(updated)


# --------------------------------------------------------------------------
# Tools — modules
# --------------------------------------------------------------------------


@mcp.tool()
async def list_modules(course_id: int) -> list[dict]:
    """List a course's modules and their items (pages, assignments, quizzes, files, links)."""
    modules = await _canvas_paginated(f"/courses/{course_id}/modules", params={"include[]": ["items"]})
    return [_slim_module(m) for m in modules]


@mcp.tool()
async def create_module(course_id: int, name: str, position: int | None = None, unlock_at: str | None = None) -> dict:
    """Create a new (unpublished) module in a course.

    - position: 1-based position in the module list; omit to append at the end.
    - unlock_at: ISO 8601 timestamp when the module unlocks for students.
    """
    module: dict = {"name": name}
    if position is not None:
        module["position"] = position
    if unlock_at:
        module["unlock_at"] = unlock_at
    created = await _canvas("POST", f"/courses/{course_id}/modules", body={"module": module})
    return _slim_module(created)


@mcp.tool()
async def update_module(
    course_id: int,
    module_id: int,
    name: str | None = None,
    published: bool | None = None,
    position: int | None = None,
) -> dict:
    """Rename, reorder, or publish/unpublish a module. Only the fields you pass are changed."""
    module: dict = {}
    if name is not None:
        module["name"] = name
    if published is not None:
        module["published"] = published
    if position is not None:
        module["position"] = position
    if not module:
        raise ValueError("Pass at least one field to update.")
    updated = await _canvas("PUT", f"/courses/{course_id}/modules/{module_id}", body={"module": module})
    return _slim_module(updated)


@mcp.tool()
async def add_module_item(
    course_id: int,
    module_id: int,
    item_type: str,
    title: str,
    content_id: int | None = None,
    page_url: str | None = None,
    external_url: str | None = None,
    position: int | None = None,
    indent: int = 0,
) -> dict:
    """Add an item to a module.

    - item_type: "Assignment", "Quiz", "File", "Discussion", "Page", "SubHeader",
      "ExternalUrl", or "ExternalTool".
    - content_id: required for Assignment/Quiz/File/Discussion (the object's Canvas id).
    - page_url: required for Page items (the page's URL slug).
    - external_url: required for ExternalUrl items.
    """
    item: dict = {"type": item_type, "title": title, "indent": indent}
    if content_id is not None:
        item["content_id"] = content_id
    if page_url:
        item["page_url"] = page_url
    if external_url:
        item["external_url"] = external_url
    if position is not None:
        item["position"] = position
    created = await _canvas("POST", f"/courses/{course_id}/modules/{module_id}/items", body={"module_item": item})
    return {
        "id": created.get("id"),
        "title": created.get("title"),
        "type": created.get("type"),
        "position": created.get("position"),
        "published": created.get("published"),
    }


@mcp.tool()
async def read_module_item(course_id: int, module_id: int, item_id: int) -> str:
    """Read the content behind one module item (get item ids from list_modules).

    Pages return their body text; File items are downloaded and their text
    extracted (PDF, DOCX, PPTX, HTML, plain text); Assignments, Quizzes, and
    Discussions return their descriptions; ExternalUrl items return the URL.
    """
    item = await _canvas("GET", f"/courses/{course_id}/modules/{module_id}/items/{item_id}")
    itype = item.get("type")
    header = f"[{itype}] {item.get('title') or ''}".rstrip()

    if itype == "Page":
        page = await _canvas("GET", f"/courses/{course_id}/pages/{item.get('page_url')}")
        return f"{header}\n\n{_strip_html(page.get('body') or '') or '[empty page]'}"
    if itype == "File":
        f = await _canvas("GET", f"/courses/{course_id}/files/{item.get('content_id')}")
        return f"{header}\n\n{await _file_text(f)}"
    if itype == "Assignment":
        a = await _canvas("GET", f"/courses/{course_id}/assignments/{item.get('content_id')}")
        meta = f"points: {a.get('points_possible')}  due: {a.get('due_at')}  published: {a.get('published')}"
        return f"{header}\n{meta}\n\n{_strip_html(a.get('description') or '') or '[no description]'}"
    if itype == "Quiz":
        q = await _canvas("GET", f"/courses/{course_id}/quizzes/{item.get('content_id')}")
        meta = f"questions: {q.get('question_count')}  points: {q.get('points_possible')}  due: {q.get('due_at')}"
        return f"{header}\n{meta}\n\n{_strip_html(q.get('description') or '') or '[no description]'}"
    if itype == "Discussion":
        d = await _canvas("GET", f"/courses/{course_id}/discussion_topics/{item.get('content_id')}")
        return f"{header}\n\n{_strip_html(d.get('message') or '') or '[no message]'}"
    if itype in ("ExternalUrl", "ExternalTool"):
        return f"{header}\nURL: {item.get('external_url')}"
    return f"{header}\n[{itype} items have no readable content]"


@mcp.tool()
async def update_module_item(
    course_id: int,
    module_id: int,
    item_id: int,
    title: str | None = None,
    position: int | None = None,
    indent: int | None = None,
    published: bool | None = None,
) -> dict:
    """Rename, reorder, re-indent, or publish/unpublish one module item."""
    item: dict = {}
    if title is not None:
        item["title"] = title
    if position is not None:
        item["position"] = position
    if indent is not None:
        item["indent"] = indent
    if published is not None:
        item["published"] = published
    if not item:
        raise ValueError("Pass at least one field to update.")
    updated = await _canvas(
        "PUT", f"/courses/{course_id}/modules/{module_id}/items/{item_id}", body={"module_item": item}
    )
    return {
        "id": updated.get("id"),
        "title": updated.get("title"),
        "type": updated.get("type"),
        "position": updated.get("position"),
        "published": updated.get("published"),
    }


@mcp.tool()
async def delete_module_item(course_id: int, module_id: int, item_id: int) -> dict:
    """Remove one item from a module. The underlying content (page, file,
    assignment...) is NOT deleted — only its entry in the module."""
    await _canvas("DELETE", f"/courses/{course_id}/modules/{module_id}/items/{item_id}")
    return {"deleted": True, "module_item_id": item_id}


@mcp.tool()
async def delete_module(course_id: int, module_id: int) -> dict:
    """Delete a module (its items are removed from the module but the underlying
    content — assignments, pages, files — is NOT deleted). This cannot be undone;
    confirm with the instructor before calling."""
    await _canvas("DELETE", f"/courses/{course_id}/modules/{module_id}")
    return {"deleted": True, "module_id": module_id}


# --------------------------------------------------------------------------
# Tools — wiki pages
# --------------------------------------------------------------------------


@mcp.tool()
async def list_pages(course_id: int, search_term: str | None = None) -> list[dict]:
    """List wiki pages in a course (page_id, url slug, title, published state)."""
    params = {"search_term": search_term} if search_term else None
    pages = await _canvas_paginated(f"/courses/{course_id}/pages", params=params)
    return [_slim_page(p) for p in pages]


@mcp.tool()
async def get_page(course_id: int, page_url_or_id: str, raw_html: bool = False) -> str:
    """Read one page's content. Pass the page's URL slug (from list_pages) or numeric page_id.

    Returns the body as plain text by default; set raw_html=True to get the
    original HTML (do that before editing a page with update_page).
    """
    page = await _canvas("GET", f"/courses/{course_id}/pages/{page_url_or_id}")
    body = page.get("body") or ""
    text = body if raw_html else _strip_html(body)
    return (
        f"# {page.get('title')} (url: {page.get('url')}, published: {page.get('published')})\n\n"
        f"{text or '[empty page]'}"
    )


@mcp.tool()
async def create_page(
    course_id: int,
    title: str,
    body_html: str,
    published: bool = False,
    module_id: int | None = None,
) -> dict:
    """Create a wiki page (lecture notes, weekly overview, syllabus, ...).

    - body_html: the page content, HTML allowed.
    - published: leave False to create a draft the instructor can review first.
    - module_id: optionally add the new page straight to that module.
    """
    page = {"title": title, "body": body_html, "published": published}
    created = await _canvas("POST", f"/courses/{course_id}/pages", body={"wiki_page": page})
    out: dict = {"page": _slim_page(created)}
    if module_id is not None and created.get("url"):
        out["module_item"] = await add_module_item(
            course_id=course_id, module_id=module_id, item_type="Page",
            title=title, page_url=created["url"],
        )
    return out


@mcp.tool()
async def update_page(
    course_id: int,
    page_url_or_id: str,
    title: str | None = None,
    body_html: str | None = None,
    published: bool | None = None,
) -> dict:
    """Update a page's title, body, or published state. Only the fields you pass are changed.

    Fetch the current HTML first with get_page(raw_html=True) when editing the body,
    so existing formatting is preserved.
    """
    page: dict = {}
    if title is not None:
        page["title"] = title
    if body_html is not None:
        page["body"] = body_html
    if published is not None:
        page["published"] = published
    if not page:
        raise ValueError("Pass at least one field to update.")
    updated = await _canvas("PUT", f"/courses/{course_id}/pages/{page_url_or_id}", body={"wiki_page": page})
    return _slim_page(updated)


# --------------------------------------------------------------------------
# Tools — course files
# --------------------------------------------------------------------------


@mcp.tool()
async def list_files(course_id: int, search_term: str | None = None) -> list[dict]:
    """List files in the course Files area (id, name, type, size). search_term filters by name."""
    params = {"search_term": search_term} if search_term else None
    files = await _canvas_paginated(f"/courses/{course_id}/files", params=params)
    return [_slim_file(f) for f in files]


@mcp.tool()
async def get_file_text(file_id: int) -> str:
    """Download a Canvas file and return its extracted text.

    Supports PDF, DOCX, PPTX, HTML, and plain-text files (txt/md/csv/...).
    Legacy binary formats (.doc, .ppt) and media files cannot be extracted.
    Find file ids with list_files or list_modules.
    """
    f = await _canvas("GET", f"/files/{file_id}")
    name = f.get("display_name") or f.get("filename") or f"file {file_id}"
    return f"=== {name} ===\n{await _file_text(f)}"


@mcp.tool()
async def upload_file(
    course_id: int,
    file_path: str,
    display_name: str | None = None,
    folder_path: str | None = None,
    module_id: int | None = None,
    on_duplicate: str = "rename",
) -> dict:
    """Upload a local file (PDF, PPTX, DOCX, images, ...) to the course Files area.

    - file_path: absolute path of the file on the machine running this server.
    - folder_path: Canvas folder like "course files/week 3"; created if missing.
    - module_id: optionally add the uploaded file straight to that module.
    - on_duplicate: "rename" (default) or "overwrite" when a same-named file exists.
    """
    path = Path(file_path)
    if not path.is_file():
        raise ValueError(f"No such file on this machine: {file_path}")
    data = path.read_bytes()
    if len(data) > MAX_ATTACHMENT_BYTES:
        raise ValueError(f"File is {len(data)} bytes, over the {MAX_ATTACHMENT_BYTES} byte limit.")
    name = display_name or path.name
    ctype = mimetypes.guess_type(name)[0] or "application/octet-stream"

    # Step 1: tell Canvas about the file; it returns a one-time upload URL.
    pre: dict = {"name": name, "size": len(data), "content_type": ctype, "on_duplicate": on_duplicate}
    if folder_path:
        pre["parent_folder_path"] = folder_path
    step1 = await _canvas("POST", f"/courses/{course_id}/files", body=pre)

    # Step 2: send the bytes to the upload URL (no Canvas auth header on this one).
    resp = await _client().post(
        step1["upload_url"],
        data=step1.get("upload_params") or {},
        files={"file": (name, data, ctype)},
        follow_redirects=False,
    )
    if resp.status_code >= 400:
        raise RuntimeError(f"File upload failed: HTTP {resp.status_code}: {resp.text[:300]}")

    # Step 3: confirm the upload (redirect or JSON "location") to get the file record.
    file_json: dict = {}
    location = resp.headers.get("Location")
    if not location:
        try:
            file_json = resp.json()
        except Exception:
            file_json = {}
        if "id" not in file_json:
            location = file_json.get("location")
    if location:
        confirm = await _throttled("GET", location, headers=_headers())
        if confirm.status_code >= 400:
            raise RuntimeError(f"Upload confirmation failed: HTTP {confirm.status_code}")
        file_json = confirm.json()
    if "attachment" in file_json:  # some Canvas instances nest the file record
        file_json = file_json["attachment"]

    out: dict = {"file": _slim_file(file_json)}
    if module_id is not None and file_json.get("id"):
        out["module_item"] = await add_module_item(
            course_id=course_id, module_id=module_id, item_type="File",
            title=name, content_id=file_json["id"],
        )
    return out


# --------------------------------------------------------------------------
# Tools — students, announcements, discussions
# --------------------------------------------------------------------------


@mcp.tool()
async def list_students(course_id: int, search_term: str | None = None) -> list[dict]:
    """List students enrolled in a course (id, name, email where visible)."""
    params: dict = {"enrollment_type[]": "student"}
    if search_term:
        params["search_term"] = search_term
    users = await _canvas_paginated(f"/courses/{course_id}/users", params=params)
    return [
        {
            "id": u.get("id"),
            "name": u.get("name"),
            "sortable_name": u.get("sortable_name"),
            "email": u.get("email") or u.get("login_id"),
        }
        for u in users
    ]


@mcp.tool()
async def create_announcement(
    course_id: int, title: str, message_html: str, delayed_post_at: str | None = None
) -> dict:
    """Post a course announcement. It is visible to students IMMEDIATELY unless
    delayed_post_at (ISO 8601) is set — confirm the text with the instructor first."""
    body: dict = {"title": title, "message": message_html, "is_announcement": True, "published": True}
    if delayed_post_at:
        body["delayed_post_at"] = delayed_post_at
    created = await _canvas("POST", f"/courses/{course_id}/discussion_topics", body=body)
    return {
        "id": created.get("id"),
        "title": created.get("title"),
        "posted_at": created.get("posted_at"),
        "delayed_post_at": created.get("delayed_post_at"),
        "html_url": created.get("html_url"),
    }


@mcp.tool()
async def list_announcements(course_id: int) -> list[dict]:
    """List the course's announcements."""
    topics = await _canvas_paginated(
        f"/courses/{course_id}/discussion_topics", params={"only_announcements": True}
    )
    return [
        {"id": t.get("id"), "title": t.get("title"), "posted_at": t.get("posted_at"), "html_url": t.get("html_url")}
        for t in topics
    ]


@mcp.tool()
async def create_discussion(
    course_id: int, title: str, message_html: str, published: bool = False, threaded: bool = True
) -> dict:
    """Create a discussion topic (unpublished draft by default).

    Add it to a module afterwards with add_module_item(item_type="Discussion", content_id=...).
    """
    body = {
        "title": title,
        "message": message_html,
        "published": published,
        "discussion_type": "threaded" if threaded else "side_comment",
    }
    created = await _canvas("POST", f"/courses/{course_id}/discussion_topics", body=body)
    return {
        "id": created.get("id"),
        "title": created.get("title"),
        "published": created.get("published"),
        "html_url": created.get("html_url"),
    }


# --------------------------------------------------------------------------
# MCP apps & visualization (read-only) — interactive panels for the client
#
# Tools below return {"mcp_app": {resource, title, data}, "summary": ...}.
# An MCP-Apps-aware host fetches the ui:// resource, renders it in a
# sandboxed iframe, and injects `data`; other clients just see the summary.
# None of these tools writes anything to Canvas.
# --------------------------------------------------------------------------


def _score_stats(scores: list[float]) -> dict:
    if not scores:
        return {"mean": 0, "min": 0, "max": 0}
    return {
        "mean": round(sum(scores) / len(scores), 1),
        "min": min(scores),
        "max": max(scores),
    }


async def _submission_split(course_id: int, assignment_id: int) -> tuple[list[float], int, int]:
    """(graded scores, submitted-but-ungraded count, unsubmitted count)."""
    subs = await _canvas_paginated(
        f"/courses/{course_id}/assignments/{assignment_id}/submissions"
    )
    scores: list[float] = []
    pending = 0
    missing = 0
    for s in subs:
        state = s.get("workflow_state")
        if state == "graded" and s.get("score") is not None:
            scores.append(float(s["score"]))
        elif s.get("submitted_at"):
            pending += 1
        else:
            missing += 1
    return scores, pending, missing


@mcp.tool()
async def get_course_overview(course_id: int) -> dict:
    """One read-only snapshot of a course: modules (with items), assignments,
    quizzes, pages, and announcements. Used by the course explorer app; also
    handy to orient before planning changes."""
    course = await _canvas("GET", f"/courses/{course_id}", params={"include[]": ["term", "total_students"]})
    modules = await _canvas_paginated(f"/courses/{course_id}/modules", params={"include[]": ["items"]})
    assignments = await _canvas_paginated(f"/courses/{course_id}/assignments")
    quizzes = await _canvas_paginated(f"/courses/{course_id}/quizzes")
    pages = await _canvas_paginated(f"/courses/{course_id}/pages")
    topics = await _canvas_paginated(
        f"/courses/{course_id}/discussion_topics", params={"only_announcements": True}
    )
    return {
        "course": _slim_course(course),
        "modules": [_slim_module(m) for m in modules],
        "assignments": [_slim_assignment(a) for a in assignments],
        "quizzes": [_slim_quiz(q) for q in quizzes],
        "pages": [_slim_page(p) for p in pages],
        "announcements": [
            {"id": t.get("id"), "title": t.get("title"), "posted_at": t.get("posted_at")}
            for t in topics
        ],
    }


@mcp.tool()
async def get_assignment_scores(course_id: int, assignment_id: int) -> dict:
    """Read-only grading snapshot for one assignment: the list of graded
    scores plus counts of pending (submitted, ungraded) and unsubmitted.
    No student names or ids are included."""
    assignment = await _canvas("GET", f"/courses/{course_id}/assignments/{assignment_id}")
    scores, pending, missing = await _submission_split(course_id, assignment_id)
    return {
        "assignment": assignment.get("name"),
        "points_possible": assignment.get("points_possible"),
        "scores": scores,
        "pending": pending,
        "missing": missing,
        **_score_stats(scores),
    }


@mcp.tool()
async def open_course_explorer() -> dict:
    """Open the interactive Canvas course explorer for the user: browse
    courses, modules, assignments (with grade charts), quizzes, pages and
    announcements, and push any item into the workflow context. Read-only."""
    courses = await _canvas_paginated(
        "/courses",
        params={"enrollment_state": "active", "include[]": ["term", "total_students"]},
    )
    slim = [_slim_course(c) for c in courses]
    return {
        "mcp_app": {
            "resource": "ui://canvas/explorer",
            "title": "Canvas course explorer",
            "data": {"courses": slim},
        },
        "summary": (
            f"Opened the Canvas course explorer ({len(slim)} active courses). "
            "The user can browse modules, assignments, quizzes, pages and "
            "announcements, view grade charts, and send items into the "
            "workflow context."
        ),
    }


@mcp.tool()
async def render_grade_distribution(course_id: int, assignment_id: int) -> dict:
    """Render an interactive histogram of graded scores for one assignment
    (plus pending/missing counts). Read-only; no student identities."""
    assignment = await _canvas("GET", f"/courses/{course_id}/assignments/{assignment_id}")
    points = assignment.get("points_possible") or 0
    scores, pending, missing = await _submission_split(course_id, assignment_id)
    top = max([points, *scores]) or 1
    n_bins = 8
    bins = [
        {
            "label": f"{round(top * i / n_bins)}–{round(top * (i + 1) / n_bins)}",
            "count": 0,
        }
        for i in range(n_bins)
    ]
    for s in scores:
        bins[min(n_bins - 1, int(s / top * n_bins))]["count"] += 1
    stats = _score_stats(scores)
    summary = (
        f"Grade distribution for '{assignment.get('name')}' "
        f"({len(scores)} graded of {len(scores) + pending + missing}): "
        f"mean {stats['mean']}/{points}, min {stats['min']}, max "
        f"{stats['max']}; {pending} submitted awaiting grading, "
        f"{missing} not submitted."
    )
    return {
        "mcp_app": {
            "resource": "ui://canvas/chart",
            "title": f"Grades — {assignment.get('name')}",
            "data": {
                "kind": "histogram",
                "title": f"Score distribution — {assignment.get('name')}",
                "subtitle": f"{len(scores)} graded · mean {stats['mean']} of {points} pts",
                "bins": bins,
                "workflow_text": summary,
            },
        },
        "summary": summary,
    }


@mcp.tool()
async def render_assignment_averages(course_id: int, max_assignments: int = 15) -> dict:
    """Render an interactive bar chart of the average score (as % of points
    possible) for each graded assignment in a course. Read-only."""
    assignments = await _canvas_paginated(f"/courses/{course_id}/assignments")
    items = []
    for a in assignments[:max_assignments]:
        points = a.get("points_possible") or 0
        if not points:
            continue
        scores, _, _ = await _submission_split(course_id, a["id"])
        if not scores:
            continue
        items.append(
            {
                "label": a.get("name") or f"assignment {a['id']}",
                "value": round(sum(scores) / len(scores) / points * 100, 1),
                "graded": len(scores),
            }
        )
    summary = (
        f"Average scores across {len(items)} graded assignments in course "
        f"{course_id}: "
        + (
            "; ".join(f"{i['label']}: {i['value']}%" for i in items)
            if items
            else "no graded submissions yet."
        )
    )
    return {
        "mcp_app": {
            "resource": "ui://canvas/chart",
            "title": "Assignment averages",
            "data": {
                "kind": "bars",
                "title": "Average score per assignment",
                "subtitle": f"% of points possible · course {course_id}",
                "items": items,
                "workflow_text": summary,
            },
        },
        "summary": summary,
    }


@mcp.tool()
async def render_course_progress(course_id: int, max_assignments: int = 15) -> dict:
    """Render an interactive donut of grading progress across a course's
    assignments: graded vs awaiting grading vs not submitted. Read-only."""
    assignments = await _canvas_paginated(f"/courses/{course_id}/assignments")
    graded = pending = missing = 0
    for a in assignments[:max_assignments]:
        scores, p, m = await _submission_split(course_id, a["id"])
        graded += len(scores)
        pending += p
        missing += m
    total = graded + pending + missing
    summary = (
        f"Grading progress for course {course_id} across "
        f"{min(len(assignments), max_assignments)} assignments: "
        f"{graded} graded, {pending} awaiting grading, {missing} not "
        f"submitted ({round(graded / total * 100) if total else 0}% graded)."
    )
    return {
        "mcp_app": {
            "resource": "ui://canvas/chart",
            "title": "Grading progress",
            "data": {
                "kind": "donut",
                "title": "Grading progress",
                "subtitle": f"course {course_id} · {total} submissions expected",
                "graded": graded,
                "pending": pending,
                "missing": missing,
                "workflow_text": summary,
            },
        },
        "summary": summary,
    }


# --------------------------------------------------------------------------
# Resources — read-only course views
# --------------------------------------------------------------------------


@mcp.resource("ui://canvas/explorer")
def explorer_app() -> str:
    """Interactive course explorer app (HTML, rendered by MCP-Apps hosts)."""
    import canvas_apps

    return canvas_apps.EXPLORER_HTML


@mcp.resource("ui://canvas/chart")
def chart_app() -> str:
    """Interactive chart app (HTML, rendered by MCP-Apps hosts)."""
    import canvas_apps

    return canvas_apps.CHART_HTML


@mcp.resource("canvas://courses")
async def courses_resource() -> str:
    """Active courses visible to the configured Canvas token."""
    courses = await _canvas_paginated("/courses", params={"enrollment_state": "active", "include[]": ["term"]})
    return json.dumps([_slim_course(c) for c in courses], indent=2)


@mcp.resource("canvas://courses/{course_id}/assignments")
async def assignments_resource(course_id: int) -> str:
    """Assignments in one course."""
    assignments = await _canvas_paginated(f"/courses/{course_id}/assignments")
    return json.dumps([_slim_assignment(a) for a in assignments], indent=2)


@mcp.resource("canvas://courses/{course_id}/quizzes")
async def quizzes_resource(course_id: int) -> str:
    """Quizzes in one course."""
    quizzes = await _canvas_paginated(f"/courses/{course_id}/quizzes")
    return json.dumps([_slim_quiz(q) for q in quizzes], indent=2)


@mcp.resource("canvas://courses/{course_id}/modules")
async def modules_resource(course_id: int) -> str:
    """Modules and module items in one course."""
    modules = await _canvas_paginated(f"/courses/{course_id}/modules", params={"include[]": ["items"]})
    return json.dumps([_slim_module(m) for m in modules], indent=2)


@mcp.resource("canvas://courses/{course_id}/pages")
async def pages_resource(course_id: int) -> str:
    """Wiki pages in one course."""
    pages = await _canvas_paginated(f"/courses/{course_id}/pages")
    return json.dumps([_slim_page(p) for p in pages], indent=2)


@mcp.resource("canvas://courses/{course_id}/files")
async def files_resource(course_id: int) -> str:
    """Files in one course's Files area."""
    files = await _canvas_paginated(f"/courses/{course_id}/files")
    return json.dumps([_slim_file(f) for f in files], indent=2)


# --------------------------------------------------------------------------
# Prompts — reusable grading / quiz workflows for the client LLM
# --------------------------------------------------------------------------


@mcp.prompt()
def grade_homework(course_id: str, assignment_id: str, rubric: str = "") -> str:
    """Grade every ungraded submission for an assignment against a rubric."""
    rubric_block = rubric.strip() or (
        "No rubric was provided. First fetch the assignment description with "
        "list_assignments and derive a sensible point breakdown from it, then "
        "show me the rubric you plan to use before grading anything."
    )
    return f"""You are grading homework in Canvas course {course_id}, assignment {assignment_id}.

Rubric / grading criteria:
{rubric_block}

Work through this process:
1. Call list_submissions(course_id={course_id}, assignment_id={assignment_id}, only_ungraded=True).
2. For each student, call get_submission_text to read their work (text entries and
   PDF/DOCX uploads are extracted automatically).
3. Assess each submission against the rubric. For each one, produce: the proposed
   score, 2-4 sentences of specific feedback citing the student's own work, and any
   red flags (empty file, off-topic, suspected plagiarism).
4. Show me the full grade table for approval BEFORE posting anything.
5. Only after I approve, call grade_submission for each student with the grade and
   the feedback text as the comment."""


@mcp.prompt()
def assess_single_submission(course_id: str, assignment_id: str, user_id: str, rubric: str = "") -> str:
    """Assess one student's submission in detail without posting a grade."""
    return f"""Fetch the submission with get_submission_text(course_id={course_id},
assignment_id={assignment_id}, user_id={user_id}) and assess it.

Rubric / criteria:
{rubric.strip() or "Derive reasonable criteria from the assignment description (see list_assignments)."}

Give me: a per-criterion breakdown with points, overall proposed grade, strengths,
weaknesses, and a draft feedback comment written directly to the student. Do NOT
call grade_submission — I will decide whether to post it."""


@mcp.prompt()
def build_quiz(course_id: str, topic: str, num_questions: str = "10", difficulty: str = "mixed") -> str:
    """Draft and create a quiz on a topic in a Canvas course."""
    return f"""Create a quiz in Canvas course {course_id} on: {topic}

Requirements: {num_questions} questions, {difficulty} difficulty, mostly multiple
choice with 4 options each, plus 1-2 short answer or essay questions if appropriate.

Process:
1. Draft all questions with answers first and show them to me for review.
2. After I approve, call create_quiz (leave it unpublished), then add_quiz_question
   for each question, marking the correct answers.
3. Report the quiz id and remind me to run publish_quiz when I'm ready to release it."""


@mcp.prompt()
def build_course_module(course_id: str, module_name: str, materials: str = "") -> str:
    """Assemble a complete course module: overview page, files, assignment, and quiz."""
    return f"""Build a complete module named "{module_name}" in Canvas course {course_id}.

Materials / notes from the instructor:
{materials.strip() or "None provided — propose content from the module name and the existing course (check list_modules and list_pages for the course's style and topics)."}

Process:
1. Look at list_modules(course_id={course_id}) to match the course's existing naming
   and structure, then show me an outline for the new module: an overview page, any
   files to upload (ask me for local paths), one assignment (with description,
   deadline, and points), and a short quiz.
2. After I approve the outline: create_module, then create_page for the overview
   (pass module_id to attach it), upload_file for local materials (also with
   module_id), create_assignment and create_quiz + add_quiz_question (both take
   module_id too) — everything as unpublished drafts.
3. Give me a summary table of what was created (ids and titles) and remind me to
   publish the module and its items when ready."""


# --------------------------------------------------------------------------


if __name__ == "__main__":
    import sys

    # stdout carries MCP protocol messages over stdio — logs must go to stderr
    # (see stdio_instruction.md).
    if not (os.environ.get("CANVAS_BASE_URL") and os.environ.get("CANVAS_API_TOKEN")):
        print("WARNING: CANVAS_BASE_URL / CANVAS_API_TOKEN not set — tool calls will fail "
              "until you configure them in .env and restart.", file=sys.stderr)
    print("Canvas MCP server starting (stdio transport)", file=sys.stderr)
    if _AUDIT_ENABLED:
        print(f"[canvas] tool-call audit log: {_AUDIT_PATH}", file=sys.stderr)
    else:
        print("[canvas] tool-call audit log DISABLED (CANVAS_MCP_AUDIT=0)", file=sys.stderr)
    # mcp.run(transport="streamable-http")  # HTTP alternative: http://{MCP_HOST}:{MCP_PORT}/mcp
    mcp.run(transport="stdio")
